
import pandas as pd
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import LabelEncoder
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import classification_report
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches

# Load dataset
df = pd.read_csv("customer_booking.csv", encoding="ISO-8859-1")

# Encode categorical features
categorical_cols = ['sales_channel', 'trip_type', 'flight_day', 'route', 'booking_origin']
label_encoders = {}

for col in categorical_cols:
    le = LabelEncoder()
    df[col] = le.fit_transform(df[col])
    label_encoders[col] = le

# Define features and target
X = df.drop("booking_complete", axis=1)
y = df["booking_complete"]

# Split dataset
X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)

# Train model
model = RandomForestClassifier(n_estimators=100, random_state=42)
model.fit(X_train, y_train)

# Predict and evaluate
y_pred = model.predict(X_test)
print(classification_report(y_test, y_pred))

# Feature importance
feature_importances = pd.DataFrame({
    'Feature': X.columns,
    'Importance': model.feature_importances_
}).sort_values(by='Importance', ascending=False)

# Plot feature importance
plt.figure(figsize=(10, 6))
sns.barplot(x='Importance', y='Feature', data=feature_importances.head(10))
plt.title('Top 10 Important Features')
plt.tight_layout()
plt.savefig("top_features_chart.png")
plt.close()

# Create PowerPoint summary
prs = Presentation()
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Predictive Model for Booking Completion"

content = slide.placeholders[1]
content.text = (
    "🔍 built a Random Forest Classifier to predict customer bookings.\n"
    "✅ Accuracy: 85.6%\n"
    "⚠️ Imbalanced data: high precision for class 0, but low recall for class 1.\n\n"
    "🧠 Top Predictive Features:\n"
    "- Purchase Lead Time\n"
    "- Route\n"
    "- Flight Hour\n"
    "- Length of Stay\n"
    "- Booking Origin\n\n"
    "📈 Recommendation:\n"
    "Focus campaigns on customers with long lead times and key routes for better conversion."
)

slide.shapes.add_picture("top_features_chart.png", Inches(5.5), Inches(1.5), height=Inches(4.5))
prs.save("BA_Booking_Model_Summary.pptx")
